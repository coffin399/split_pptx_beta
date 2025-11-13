#!/usr/bin/env python3
"""GUI application to generate script-style slides from PPTX notes."""

from __future__ import annotations

import logging
import os
import platform
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Callable, Iterable, List, Optional, Tuple


@dataclass
class Segment:
    text: str
    speaker: Optional[str] = None

# Import disk cache for memory optimization
try:
    from cache import get_cache, cleanup_cache
    CACHE_AVAILABLE = True
except ImportError:
    CACHE_AVAILABLE = False
    def get_cache():
        return None
    def cleanup_cache():
        pass

try:  # GUI dependencies are optional in headless deployments (e.g., Render.com)
    from PySide6.QtCore import Qt, QThread, Signal
    from PySide6.QtGui import QAction, QFont
    from PySide6.QtWidgets import (
        QApplication,
        QFileDialog,
        QLabel,
        QLineEdit,
        QMainWindow,
        QMenuBar,
        QMessageBox,
        QPushButton,
        QStatusBar,
        QTextEdit,
        QToolBar,
        QVBoxLayout,
        QWidget,
        QHBoxLayout,
    )
    PYSIDE_AVAILABLE = True
except Exception as exc:  # noqa: BLE001 - broad to catch missing libGL, etc.
    PYSIDE_AVAILABLE = False
    PYSIDE_IMPORT_ERROR = exc
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)

LOGGER = logging.getLogger("pptx_script")

MAX_CHARS_PER_SLIDE = 200
OUTPUT_FILENAME = "スクリプトスライド_自動生成.pptx"
FONT_NAME = "メイリオ"
FONT_SIZE_PT = 40
FONT_BOLD = True
PAGE_INDICATOR_COLOR = RGBColor(0x00, 0xB0, 0xF0)
DEFAULT_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TEXTBOX_POSITION = {
    "left": Cm(0.79),
    "top": Cm(0.8),
    "width": Cm(22.0),  # Reduced to avoid overlap with thumbnail
    "height": Cm(15.6),
}
THUMBNAIL_WIDTH_CM = 8.0
THUMBNAIL_MARGIN_CM = 0.5
DEFAULT_THUMBNAIL_DPI = 150  # Reduced from higher values to save memory
LOW_MEMORY_DPI = 100  # Further reduced for low memory situations
ULTRA_LOW_DPI = 80   # For very large presentations (100+ slides)
EMU_PER_INCH = 914400
SPEAKER_PATTERN = re.compile(r"^\s*(話者\d+)[:：]\s*(.*)$")

SPEAKER_COLORS = {
    "話者1": RGBColor(0xFF, 0xFF, 0x00),
    "話者2": RGBColor(0x00, 0xFF, 0xFF),
    "話者3": RGBColor(0x00, 0xF9, 0x00),
}

FONT_PATH_CANDIDATES = [
    "C:/Windows/Fonts/meiryo.ttc",
    "C:/Windows/Fonts/meiryo.ttf",
    "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc",
    "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
    "/Library/Fonts/Osaka.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]

_FONT_CACHE: dict[int, ImageFont.ImageFont] = {}


def clear_font_cache() -> None:
    """Clear font cache to free memory."""
    global _FONT_CACHE
    for font_size, font in _FONT_CACHE.items():
        try:
            # Font objects don't have explicit close method, but we can clear references
            del font
        except Exception:
            pass
    _FONT_CACHE.clear()


def get_optimal_dpi(slide_count: int) -> int:
    """Get optimal DPI based on slide count to balance quality and memory usage."""
    # Force a low DPI to minimize memory usage even for small decks.
    # Visual fidelity remains acceptable for thumbnails while keeping PDF conversion lightweight.
    return 50


def get_font(size_pt: float) -> ImageFont.ImageFont:
    rounded = int(round(size_pt)) if size_pt else 20
    if rounded in _FONT_CACHE:
        return _FONT_CACHE[rounded]
    for candidate in FONT_PATH_CANDIDATES:
        if Path(candidate).exists():
            try:
                font = ImageFont.truetype(candidate, rounded)
                _FONT_CACHE[rounded] = font
                return font
            except OSError:
                continue
    font = ImageFont.load_default()
    _FONT_CACHE[rounded] = font
    return font


def rgb_color_tuple(color: Optional[RGBColor], default=(0, 0, 0)) -> tuple[int, int, int]:
    if color is None:
        return default
    try:
        return (color[0], color[1], color[2])
    except (TypeError, IndexError):
        return default


def draw_text_block(
    image: Image.Image,
    shape,
    dpi: int = DEFAULT_THUMBNAIL_DPI,
) -> bool:
    if not shape.has_text_frame:
        return False
    text_frame = shape.text_frame
    draw = ImageDraw.Draw(image)
    left = emu_to_px(int(shape.left), dpi)
    top = emu_to_px(int(shape.top), dpi)
    width = emu_to_px(int(shape.width), dpi)

    paragraphs: List[str] = []
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            text = "".join(run.text for run in paragraph.runs)
        else:
            text = paragraph.text or ""
        paragraphs.append(text)

    if not any(p.strip() for p in paragraphs):
        return False

    first_run = None
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                first_run = run
                break
        if first_run:
            break
    font_size = (
        first_run.font.size.pt
        if first_run and first_run.font.size
        else 24
    )
    color = rgb_color_tuple(first_run.font.color.rgb if first_run and first_run.font.color and first_run.font.color.rgb else None)
    font = get_font(font_size)

    lines: List[str] = []
    for paragraph_text in paragraphs:
        if not paragraph_text:
            lines.append("")
            continue
        current = ""
        for char in paragraph_text:
            test = current + char
            if draw.textlength(test, font=font) <= width or not current:
                current = test
            else:
                lines.append(current)
                current = char
        lines.append(current)
    line_height = font.getbbox("あ")[3] if hasattr(font, "getbbox") else font.size
    y = top
    drew_text = False
    for line in lines:
        draw.text((left, y), line, font=font, fill=color)
        y += line_height
        if line:
            drew_text = True

    return drew_text


def segment_line(text: str, max_chars: int) -> List[str]:
    """Split a line into punctuation-aware segments within the max length."""
    # More comprehensive punctuation pattern for better line breaks
    pattern = re.compile(r"[^。、，,.！？!?；;：:]+[。、，,.！？!?；;：:]?")
    segments: List[str] = []
    for token in pattern.findall(text):
        trimmed = token.strip()
        if not trimmed:
            continue
        if len(trimmed) <= max_chars:
            segments.append(trimmed)
        else:
            # For long segments, try to break at natural points
            break_patterns = [
                re.compile(r"(.{1,80}[、,，])(.+)"),
                re.compile(r"(.{1,120}[。.！？!?])(.+)"),
                re.compile(r"(.{1,60}[^a-zA-Z0-9]{1,2})(.+)"),
            ]
            broken = False
            for break_pattern in break_patterns:
                match = break_pattern.match(trimmed)
                if match:
                    segments.extend([part.strip() for part in match.groups() if part.strip()])
                    broken = True
                    break
            if not broken:
                # Fallback: break at max_chars
                for start in range(0, len(trimmed), max_chars):
                    segments.append(trimmed[start : start + max_chars])
    if not segments and text:
        # Fallback when regex fails (e.g., single punctuation)
        for start in range(0, len(text), max_chars):
            segments.append(text[start : start + max_chars])
    return segments


def build_segments(notes_text: str, max_chars: int) -> List[Segment]:
    """Convert raw notes into ordered segments with speaker metadata."""
    segments: List[Segment] = []
    normalized = notes_text.replace("\r\n", "\n").replace("\r", "\n")
    for raw_line in normalized.split("\n"):
        stripped = raw_line.strip()
        if not stripped:
            segments.append(Segment("", None))
            continue
        match = SPEAKER_PATTERN.match(stripped)
        if match:
            speaker_label = match.group(1)
            content = match.group(2).strip()
        else:
            speaker_label = None
            content = stripped
        line_segments = segment_line(content, max_chars)
        if not line_segments:
            segments.append(Segment(content, speaker_label))
            continue
        for idx, part in enumerate(line_segments):
            if idx == 0 and speaker_label:
                display_text = f"{speaker_label}：{part}"
            else:
                display_text = part
            segments.append(Segment(display_text, speaker_label if idx == 0 else None))
    return segments


def chunk_segments(segments: Iterable[Segment], max_chars: int) -> List[List[Segment]]:
    """Group segments into chunks that respect the character limit."""
    chunks: List[List[Segment]] = []
    current: List[Segment] = []
    current_len = 0

    for seg in segments:
        seg_len = len(seg.text)
        additional = max(seg_len, 1)  # blank lines count as 1
        if current:
            additional += 1  # newline between segments
        if current and current_len + additional > max_chars:
            chunks.append(current)
            current = [seg]
            current_len = seg_len
        else:
            current.append(seg)
            if current_len == 0:
                current_len = seg_len
            else:
                current_len += additional
    if current:
        chunks.append(current)
    return chunks


def speaker_color(speaker: Optional[str]) -> RGBColor:
    return SPEAKER_COLORS.get(speaker or "", DEFAULT_FONT_COLOR)


def log(message: str, reporter: Optional[Callable[[str], None]]) -> None:
    LOGGER.info(message)
    if reporter:
        reporter(message)


def emu_to_px(value: int, dpi: int = DEFAULT_THUMBNAIL_DPI) -> int:
    return max(1, int(round(value * dpi / EMU_PER_INCH)))


def get_font(size_pt: float) -> ImageFont.ImageFont:
    rounded = int(round(size_pt)) if size_pt else 20
    if rounded in _FONT_CACHE:
        return _FONT_CACHE[rounded]
    for candidate in FONT_PATH_CANDIDATES:
        if Path(candidate).exists():
            try:
                font = ImageFont.truetype(candidate, rounded)
                _FONT_CACHE[rounded] = font
                return font
            except OSError:
                continue
    font = ImageFont.load_default()
    _FONT_CACHE[rounded] = font
    return font


def rgb_color_tuple(color: Optional[RGBColor], default=(0, 0, 0)) -> tuple[int, int, int]:
    if color is None:
        return default
    try:
        return (color[0], color[1], color[2])
    except (TypeError, IndexError):
        return default


def draw_text_block(
    image: Image.Image,
    shape,
    dpi: int = DEFAULT_THUMBNAIL_DPI,
) -> bool:
    if not shape.has_text_frame:
        return False
    text_frame = shape.text_frame
    draw = ImageDraw.Draw(image)
    left = emu_to_px(int(shape.left), dpi)
    top = emu_to_px(int(shape.top), dpi)
    width = emu_to_px(int(shape.width), dpi)

    paragraphs: List[str] = []
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            text = "".join(run.text for run in paragraph.runs)
        else:
            text = paragraph.text or ""
        paragraphs.append(text)

    if not any(p.strip() for p in paragraphs):
        return False

    first_run = None
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                first_run = run
                break
        if first_run:
            break
    font_size = (
        first_run.font.size.pt
        if first_run and first_run.font.size
        else 24
    )
    color = rgb_color_tuple(first_run.font.color.rgb if first_run and first_run.font.color and first_run.font.color.rgb else None)
    font = get_font(font_size)

    lines: List[str] = []
    for paragraph_text in paragraphs:
        if not paragraph_text:
            lines.append("")
            continue
        current = ""
        for char in paragraph_text:
            test = current + char
            if draw.textlength(test, font=font) <= width or not current:
                current = test
            else:
                lines.append(current)
                current = char
        lines.append(current)
    line_height = font.getbbox("あ")[3] if hasattr(font, "getbbox") else font.size
    y = top
    drew_text = False
    for line in lines:
        draw.text((left, y), line, font=font, fill=color)
        y += line_height
        if line:
            drew_text = True

    return drew_text


def draw_shape_fill(image: Image.Image, shape, dpi: int = DEFAULT_THUMBNAIL_DPI) -> bool:
    fill = shape.fill
    if not fill or fill.type != MSO_FILL.SOLID:
        return False
    color = rgb_color_tuple(fill.fore_color.rgb if fill.fore_color.type is not None else None, default=(255, 255, 255))
    left = emu_to_px(int(shape.left), dpi)
    top = emu_to_px(int(shape.top), dpi)
    width = emu_to_px(int(shape.width), dpi)
    height = emu_to_px(int(shape.height), dpi)
    ImageDraw.Draw(image).rectangle(
        [left, top, left + width, top + height],
        fill=color,
    )
    return True


def draw_picture(image: Image.Image, shape, dpi: int = DEFAULT_THUMBNAIL_DPI) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    try:
        blob = shape.image.blob
    except Exception:
        return False
    with Image.open(BytesIO(blob)) as pic:
        pic = pic.convert("RGBA")
        width = emu_to_px(int(shape.width), dpi)
        height = emu_to_px(int(shape.height), dpi)
        if width > 0 and height > 0:
            pic = pic.resize((width, height), Image.LANCZOS)
        left = emu_to_px(int(shape.left), dpi)
        top = emu_to_px(int(shape.top), dpi)
        image.paste(pic, (left, top), pic if pic.mode == "RGBA" else None)
    return True


def slide_background_color(slide) -> tuple[int, int, int]:
    fill = slide.background.fill
    if fill and fill.type == MSO_FILL.SOLID:
        return rgb_color_tuple(fill.fore_color.rgb if fill.fore_color and fill.fore_color.rgb else None, default=(255, 255, 255))
    return (30, 30, 30)


def _draw_placeholder_notice(image: Image.Image) -> None:
    draw = ImageDraw.Draw(image)
    overlay_color = (40, 40, 40)
    draw.rectangle([(0, 0), (image.width, image.height)], fill=overlay_color)
    title_font = get_font(48)
    body_font = get_font(28)
    lines = [
        "サムネイルを内部描画しましたが",
        "表示できる要素が見つかりませんでした",
        "LibreOffice をインストールすると正確なプレビューが作成できます",
    ]
    current_y = image.height // 2 - (len(lines) * 50) // 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=body_font)
        text_width = bbox[2] - bbox[0]
        draw.text(
            ((image.width - text_width) / 2, current_y),
            line,
            font=body_font,
            fill=(220, 220, 220),
        )
        current_y += 50


def render_slide_to_image(
    slide,
    slide_width: int,
    slide_height: int,
    output_path: Path,
    dpi: int = DEFAULT_THUMBNAIL_DPI,
) -> bool:
    width_px = emu_to_px(slide_width, dpi)
    height_px = emu_to_px(slide_height, dpi)
    background = slide_background_color(slide)
    image = Image.new("RGB", (width_px, height_px), color=background)

    drawn_any = False
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if draw_picture(image, shape, dpi):
                    drawn_any = True
            else:
                if draw_shape_fill(image, shape, dpi):
                    drawn_any = True
        except Exception:
            continue

    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            if draw_text_block(image, shape, dpi):
                drawn_any = True
        except Exception:
            continue

    if not drawn_any:
        _draw_placeholder_notice(image)

    image.save(output_path)
    return drawn_any


def ensure_blank_presentation() -> Presentation:
    prs = Presentation()
    if prs.slides:
        slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
        slide_id = slide_id_list[0]
        slide_id_list.remove(slide_id)
    return prs


def apply_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)


def add_textbox(
    slide,
    chunk: List[Segment],
    slide_width,
    slide_height,
    thumbnail_geometry: Optional[Tuple[float, float, float, float]] = None,
) -> None:
    left = TEXTBOX_POSITION["left"]
    top = TEXTBOX_POSITION["top"]
    if thumbnail_geometry:
        thumb_left, _, _, _ = thumbnail_geometry
        right_boundary = thumb_left - Cm(THUMBNAIL_MARGIN_CM)
        available_width = max(Cm(6), right_boundary - left)
    else:
        right_margin = Cm(THUMBNAIL_MARGIN_CM)
        available_width = max(Cm(10), slide_width - left - right_margin)
    available_height = max(
        Cm(5),
        slide_height - top - Cm(THUMBNAIL_MARGIN_CM + 0.5),
    )

    textbox = slide.shapes.add_textbox(
        left,
        top,
        available_width,
        available_height,
    )
    text_frame = textbox.text_frame
    text_frame.text = ""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Auto-fit text to shape
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_bottom = Pt(8)
    text_frame.margin_top = Pt(8)
    text_frame.margin_left = Pt(8)
    text_frame.margin_right = Pt(8)

    for idx, segment in enumerate(chunk):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.clear()
        run = paragraph.add_run()
        run.text = segment.text
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(12)  # Add spacing between paragraphs
        paragraph.space_before = Pt(0)
        font = run.font
        font.name = FONT_NAME
        font.size = Pt(FONT_SIZE_PT)
        font.bold = FONT_BOLD
        font.color.rgb = speaker_color(segment.speaker)


def add_page_indicator(
    slide,
    index: int,
    total: int,
    slide_width,
    slide_height,
    thumbnail_geometry: Optional[Tuple[float, float, float, float]] = None,
) -> None:
    if total <= 1:
        return
    # Position page indicator above thumbnail
    margin = Cm(0.1)
    indicator_width = Cm(4)
    indicator_height = Cm(1.5)

    if thumbnail_geometry:
        thumb_left, thumb_top, thumb_width, thumb_height = thumbnail_geometry
        indicator_left = max(margin, thumb_left - indicator_width - margin)
        indicator_top = thumb_top + max(0, (thumb_height - indicator_height) / 2)
    else:
        thumbnail_height = Cm(THUMBNAIL_WIDTH_CM * 9/16)  # Assuming 16:9 aspect ratio
        indicator_left = slide_width - indicator_width - margin
        indicator_top = slide_height - thumbnail_height - indicator_height - Cm(0.5)

    textbox = slide.shapes.add_textbox(
        indicator_left,
        indicator_top,
        indicator_width,
        indicator_height,
    )
    text_frame = textbox.text_frame
    text_frame.text = ""
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = f"{index}/{total}"
    paragraph.alignment = PP_ALIGN.RIGHT
    font = paragraph.font
    font.name = FONT_NAME
    font.size = Pt(FONT_SIZE_PT)
    font.bold = True
    font.color.rgb = PAGE_INDICATOR_COLOR


def add_thumbnail(
    slide,
    image_path: Path,
    slide_width,
    slide_height,
) -> Optional[Tuple[float, float, float, float]]:
    if not image_path.exists():
        return None
    margin = Cm(0.1)  # Minimal margin for tight positioning
    img = None
    try:
        img = Image.open(image_path)
        width_cm = THUMBNAIL_WIDTH_CM
        height_cm = width_cm * img.height / img.width
    except Exception as exc:
        log(f"サムネイル画像の読み込みに失敗しました: {exc}", None)
        return None
    finally:
        if img is not None:
            try:
                img.close()
            except Exception:
                pass
    
    width = Cm(width_cm)
    height = Cm(height_cm)
    # Position thumbnail exactly at bottom-right corner
    left = slide_width - width - margin
    top = slide_height - height - margin
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    return (left, top, width, height)


def create_placeholder_thumbnail(index: int, output_dir: Path) -> Path:
    width, height = 1600, 900
    image = None
    draw = None
    try:
        image = Image.new("RGB", (width, height), color=(30, 30, 30))
        draw = ImageDraw.Draw(image)
        title = f"スライド {index}"
        body_lines = [
            "LibreOffice/soffice が見つからないため",
            "サムネイルを簡易表示に切り替えました",
            "フルプレビューを得るには LibreOffice をインストールしてください",
        ]
        title_font = get_font(56)
        body_font = get_font(28)

        title_bbox = draw.textbbox((0, 0), title, font=title_font)
        title_width = title_bbox[2] - title_bbox[0]
        draw.text(
            ((width - title_width) / 2, height * 0.3),
            title,
            font=title_font,
            fill=(240, 240, 240),
        )

        start_y = height * 0.5
        line_spacing = 45
        for idx, line in enumerate(body_lines):
            bbox = draw.textbbox((0, 0), line, font=body_font)
            text_width = bbox[2] - bbox[0]
            draw.text(
                ((width - text_width) / 2, start_y + idx * line_spacing),
                line,
                font=body_font,
                fill=(200, 200, 200),
            )
        placeholder_path = output_dir / f"placeholder_slide_{index}.png"
        image.save(placeholder_path)
        return placeholder_path
    finally:
        # Clean up resources
        if draw is not None:
            del draw
        if image is not None:
            try:
                image.close()
            except Exception:
                pass


def _normalize_powerpoint_export(png_dir: Path) -> List[Path]:
    results: List[Path] = []
    if not png_dir.exists():
        return results
    for pattern in ("Slide*.PNG", "Slide*.png"):
        results.extend(sorted(png_dir.glob(pattern)))
    return sorted(results)


def _export_with_powerpoint_windows(
    pptx_path: Path,
    destination: Path,
    reporter: Optional[Callable[[str], None]],
) -> List[Path]:
    try:
        import win32com.client  # type: ignore[import]
    except ImportError:
        log("PowerPoint連携には pywin32 が必要です。'pip install pywin32' を実行してください。", reporter)
        return []

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    except Exception as exc:  # noqa: BLE001
        log(f"PowerPoint の起動に失敗しました: {exc}", reporter)
        return []

    destination.mkdir(parents=True, exist_ok=True)
    presentation = None
    exported: List[Path] = []
    try:
        presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        presentation.Export(str(destination), "PNG")
        exported = _normalize_powerpoint_export(destination)
    except Exception as exc:  # noqa: BLE001
        log(f"PowerPoint からのサムネイル書き出しに失敗しました: {exc}", reporter)
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:  # noqa: BLE001
                pass
        try:
            powerpoint.Quit()
        except Exception:  # noqa: BLE001
            pass
    return exported


def _export_with_powerpoint_macos(
    pptx_path: Path,
    destination: Path,
    reporter: Optional[Callable[[str], None]],
) -> List[Path]:
    destination.mkdir(parents=True, exist_ok=True)
    pptx_posix = pptx_path.as_posix().replace("\"", "\\\"")
    dest_posix = destination.as_posix().replace("\"", "\\\"")
    
    # Use proper AppleScript syntax for PowerPoint export
    applescript = f'''
    tell application "Microsoft PowerPoint"
        activate
        set thePresentation to open POSIX file "{pptx_posix}" with read only
        save thePresentation in POSIX file "{dest_posix}" as save as picture file format PNG
        close thePresentation saving no
    end tell
    '''
    
    # Clean up the applescript for proper formatting
    applescript = applescript.strip()
    
    try:
        result = subprocess.run([
            "osascript",
            "-e",
            applescript,
        ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        print(f"AppleScript execution completed: {result.returncode}")
    except FileNotFoundError:
        log("osascript コマンドが見つかりませんでした。macOS でのみ利用できます。", reporter)
        return []
    except subprocess.CalledProcessError as exc:
        error_msg = exc.stderr.strip() if exc.stderr else str(exc)
        log(
            f"PowerPoint (macOS) からのサムネイル書き出しに失敗しました: {error_msg}",
            reporter,
        )
        print(f"AppleScript failed with return code {exc.returncode}")
        print(f"stderr: {exc.stderr}")
        print(f"stdout: {exc.stdout}")
        print(f"Script that failed: {repr(applescript)}")
        return []
    exported = _normalize_powerpoint_export(destination)
    if not exported:
        log("PowerPoint から書き出された PNG が見つかりませんでした。", reporter)
    return exported


def _export_thumbnails_via_pdf(
    pptx_path: Path,
    slide_count: int,
    persistent_dir: Path,
    reporter: Optional[Callable[[str], None]],
) -> List[Path]:
    """Export thumbnails via PDF conversion with memory-optimized DPI."""
    pdf_workspace = persistent_dir / "pdf_workspace"
    pdf_workspace.mkdir(parents=True, exist_ok=True)
    pdf_path = pdf_workspace / pptx_path.with_suffix(".pdf").name

    # Use optimal DPI based on slide count
    optimal_dpi = get_optimal_dpi(slide_count)
    log(f"スライド数: {slide_count}, DPI: {optimal_dpi} (メモリ最適化)", reporter)

    soffice = next((cmd for cmd in ("soffice", "libreoffice") if shutil.which(cmd)), None)
    if not soffice:
        log(
            "LibreOffice/soffice が見つかりません。PDFベースのサムネイル生成をスキップします。",
            reporter,
        )
        return []

    try:
        env = os.environ.copy()
        # Limit memory usage for LibreOffice
        env["LIBREOFFICE_USE_SYSTEM_LIBS"] = "1"
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(pdf_workspace),
                str(pptx_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            env=env,
        )
    except (subprocess.CalledProcessError, FileNotFoundError) as exc:
        log(
            f"LibreOffice での PDF 変換に失敗しました: {exc}. 別の方法にフォールバックします。",
            reporter,
        )
        return []

    if not pdf_path.exists():
        pdf_candidates = sorted(pdf_workspace.glob("*.pdf"))
        if not pdf_candidates:
            log("PDF 変換結果が見つかりませんでした。別の方法にフォールバックします。", reporter)
            return []
        pdf_path = pdf_candidates[0]

    poppler_path = os.getenv("POPPLER_PATH") or None
    
    # Update font cache for better Japanese font detection
    try:
        subprocess.run(["fc-cache", "-fv"], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except (subprocess.CalledProcessError, FileNotFoundError):
        pass  # Font cache update is optional
    
    image_output_dir = persistent_dir / "pdf_images"
    image_output_dir.mkdir(parents=True, exist_ok=True)

    exports: List[Path] = []
    batch_size = 15 if slide_count > 120 else 25 if slide_count > 60 else 40

    for batch_start in range(1, slide_count + 1, batch_size):
        batch_end = min(batch_start + batch_size - 1, slide_count)
        log(
            f"PDF バッチ変換開始: ページ {batch_start}-{batch_end} (全 {slide_count} ページ)",
            reporter,
        )
        try:
            batch_paths = convert_from_path(
                str(pdf_path),
                dpi=optimal_dpi,
                poppler_path=poppler_path,
                fmt="png",
                output_folder=str(image_output_dir),
                output_file=f"slide_{batch_start:03d}",
                paths_only=True,
                first_page=batch_start,
                last_page=batch_end,
                thread_count=1,
                use_pdftocairo=True,
            )
        except Exception as exc:
            log(
                f"PDF からの画像変換 (ページ {batch_start}-{batch_end}) に失敗しました: {exc}。別の方法にフォールバックします。",
                reporter,
            )
            shutil.rmtree(image_output_dir, ignore_errors=True)
            return []

        expected_count = batch_end - batch_start + 1
        if len(batch_paths) != expected_count:
            log(
                f"PDF から生成されたページ数 ({len(batch_paths)}) が期待値 ({expected_count}) と一致しません (ページ {batch_start}-{batch_end})。",
                reporter,
            )
            shutil.rmtree(image_output_dir, ignore_errors=True)
            return []

        for offset, image_path in enumerate(sorted(batch_paths), start=0):
            slide_index = batch_start + offset
            src = Path(image_path)
            dest = persistent_dir / f"slide_{slide_index:03d}.png"
            try:
                shutil.move(str(src), dest)
            except Exception as exc:
                log(f"画像 {slide_index} の保存に失敗しました: {exc}", reporter)
                continue
            exports.append(dest)

        log(
            f"PDF バッチ変換完了: ページ {batch_start}-{batch_end} -> {len(batch_paths)} 枚", 
            reporter,
        )

    shutil.rmtree(image_output_dir, ignore_errors=True)

    try:
        pdf_path.unlink(missing_ok=True)
    except Exception:
        pass

    return exports


def generate_thumbnails(
    prs: Presentation,
    pptx_path: Path,
    reporter: Optional[Callable[[str], None]],
) -> List[Optional[Path]]:
    slide_count = len(prs.slides)
    thumbnails: List[Optional[Path]] = [None] * slide_count
    optimal_dpi = get_optimal_dpi(slide_count)
    
    # Try to use disk cache for memory efficiency
    cache = get_cache() if CACHE_AVAILABLE else None
    if cache:
        cache_stats = cache.get_stats()
        log(f"キャッシュ使用量: {cache_stats['total_size_mb']}MB ({cache_stats['usage_percent']}%)", reporter)
        
        # Check cache for existing thumbnails
        cache_hits = 0
        for idx in range(slide_count):
            cached_path = cache.get(pptx_path, idx, optimal_dpi)
            if cached_path:
                thumbnails[idx] = cached_path
                cache_hits += 1
        
        if cache_hits > 0:
            log(f"キャッシュから {cache_hits}/{slide_count} 枚のサムネイルを読み込みました", reporter)

            # Only generate missing thumbnails
            missing_indices = [i for i, path in enumerate(thumbnails) if path is None]
            if not missing_indices:
                log("すべてのサムネイルがキャッシュに存在しました", reporter)
                return thumbnails, None
    
    persistent_dir = Path(tempfile.mkdtemp(prefix="pptx_thumbs_"))

    pdf_exports = _export_thumbnails_via_pdf(pptx_path, slide_count, persistent_dir, reporter)
    if pdf_exports:
        log("PDF を経由したサムネイル生成に成功しました。", reporter)
        
        # Store in cache if available
        if cache:
            for idx, path in enumerate(pdf_exports):
                cache.put(pptx_path, idx, optimal_dpi, path)
            log(f"{len(pdf_exports)} 枚のサムネイルをキャッシュに保存しました", reporter)
        
        return pdf_exports, persistent_dir

    system = platform.system()
    external_dir = persistent_dir / "external"
    external_exports: List[Path] = []
    if system == "Windows":
        external_exports = _export_with_powerpoint_windows(pptx_path, external_dir, reporter)
    elif system == "Darwin":
        external_exports = _export_with_powerpoint_macos(pptx_path, external_dir, reporter)

    if external_exports:
        if len(external_exports) == slide_count:
            log("PowerPoint を使用してサムネイルを取得しました。", reporter)
            for idx, path in enumerate(external_exports):
                dest = persistent_dir / f"slide_{idx + 1:03d}.png"
                shutil.copy2(path, dest)
                thumbnails[idx] = dest
        else:
            log(
                "PowerPoint から取得したサムネイル数がスライド数と一致しません。内部レンダリングにフォールバックします。",
                reporter,
            )

    soffice = next((cmd for cmd in ("soffice", "libreoffice") if shutil.which(cmd)), None)
    if soffice:
        with tempfile.TemporaryDirectory(prefix="pptx_lo_") as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            try:
                subprocess.run(
                    [
                        soffice,
                        "--headless",
                        "--convert-to",
                        "png",
                        "--outdir",
                        str(tmp_dir),
                        str(pptx_path),
                    ],
                    check=True,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                )
                png_files = sorted(tmp_dir.glob("*.png"))
                if len(png_files) == slide_count:
                    for idx, path in enumerate(png_files):
                        dest = persistent_dir / f"slide_{idx + 1:03d}.png"
                        shutil.copy2(path, dest)
                        thumbnails[idx] = dest
                else:
                    log(
                        "LibreOffice で生成されたサムネイル数がスライド数と一致しません。内部レンダリングを使用します。",
                        reporter,
                    )
            except (subprocess.CalledProcessError, FileNotFoundError) as exc:
                log(f"サムネイル生成に失敗しました: {exc}. 内部レンダリングを使用します。", reporter)
    else:
        log(
            "LibreOffice/soffice が見つかりません。内部レンダリングでサムネイルを生成します。",
            reporter,
        )

    for idx, slide in enumerate(prs.slides, start=1):
        if thumbnails[idx - 1] is not None:
            continue
        dest = persistent_dir / f"slide_{idx:03d}.png"
        try:
            rendered = render_slide_to_image(slide, prs.slide_width, prs.slide_height, dest)
            if rendered:
                thumbnails[idx - 1] = dest
            else:
                log(
                    f"スライド {idx}: 内部レンダリング結果に表示要素がありません。プレースホルダーに切り替えます。",
                    reporter,
                )
                thumbnails[idx - 1] = create_placeholder_thumbnail(idx, persistent_dir)
        except Exception as exc:
            log(
                f"スライド {idx} のレンダリングに失敗しました: {exc}. プレースホルダーに切り替えます。",
                reporter,
            )
            thumbnails[idx - 1] = create_placeholder_thumbnail(idx, persistent_dir)

    return thumbnails, persistent_dir


def cleanup_thumbnail_dir(thumbnail_dir: Optional[Path]) -> None:
    """Clean up thumbnail temporary directory to prevent memory leaks."""
    if thumbnail_dir and thumbnail_dir.exists():
        try:
            shutil.rmtree(thumbnail_dir, ignore_errors=True)
        except Exception:
            pass


def generate_script_slides(input_file: Path, output_dir: Path, reporter: Optional[Callable[[str], None]]) -> Path:
    prs = Presentation(str(input_file))
    output_prs = ensure_blank_presentation()
    output_prs.slide_width = prs.slide_width
    output_prs.slide_height = prs.slide_height
    blank_layout = output_prs.slide_layouts[6]

    thumbnails, thumbnail_dir = generate_thumbnails(prs, input_file, reporter)

    created = 0
    try:
        for slide_index, slide in enumerate(prs.slides, start=1):
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame else ""
            if not notes.strip():
                log(f"スライド {slide_index}: ノートが空のためスキップします。", reporter)
                continue
            segments = build_segments(notes, MAX_CHARS_PER_SLIDE)
            chunks = chunk_segments(segments, MAX_CHARS_PER_SLIDE)
            log(
                f"スライド {slide_index}: {len(segments)} セグメント -> {len(chunks)} 枚に分割", reporter
            )
            thumbnail_geometry = thumbnails[slide_index - 1] if slide_index - 1 < len(thumbnails) else None

            for part_idx, chunk in enumerate(chunks, start=1):
                log(
                    f"スライド {slide_index}/{len(prs.slides)} - パート {part_idx}/{len(chunks)} を生成中", 
                    reporter,
                )
                new_slide = output_prs.slides.add_slide(blank_layout)
                apply_background(new_slide)
                thumbnail_path = thumbnails[slide_index - 1] if slide_index - 1 < len(thumbnails) else None
                thumbnail_geom = None
                if thumbnail_path:
                    thumbnail_geom = add_thumbnail(new_slide, thumbnail_path, output_prs.slide_width, output_prs.slide_height)
                add_textbox(new_slide, chunk, output_prs.slide_width, output_prs.slide_height, thumbnail_geom)
                add_page_indicator(
                    new_slide,
                    part_idx,
                    len(chunks),
                    output_prs.slide_width,
                    output_prs.slide_height,
                    thumbnail_geom,
                )
                created += 1

        if created == 0:
            raise ValueError("ノートから生成できるスライドがありませんでした。")

        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / OUTPUT_FILENAME
        output_prs.save(str(output_path))
        log(f"生成完了: {created} 枚 -> {output_path}", reporter)
        return output_path
    finally:
        # Clean up thumbnail directory to prevent memory leaks
        cleanup_thumbnail_dir(thumbnail_dir)
        # Clear font cache to free memory
        clear_font_cache()


if PYSIDE_AVAILABLE:

    class ConversionThread(QThread):
        progress = Signal(str)
        finished = Signal(bool, str)

        def __init__(self, input_file: Path, output_dir: Path):
            super().__init__()
            self.input_file = input_file
            self.output_dir = output_dir

        def run(self) -> None:  # noqa: D401
            try:
                generate_script_slides(
                    self.input_file,
                    self.output_dir,
                    self.progress.emit,
                )
                self.finished.emit(True, "処理が完了しました。")
            except Exception as exc:  # noqa: BLE001 - keep GUI responsive
                self.finished.emit(False, f"エラーが発生しました: {exc}")


    class MainWindow(QMainWindow):
        def __init__(self) -> None:
            super().__init__()
            self.setWindowTitle("ノート→スクリプトスライド変換")
            self.resize(960, 720)
            self.worker: Optional[ConversionThread] = None
            self._build_ui()

        def _build_ui(self) -> None:
            central = QWidget(self)
            layout = QVBoxLayout(central)

            form_layout = QHBoxLayout()
            layout.addLayout(form_layout)

            self.input_edit = QLineEdit(self)
            self.input_edit.setPlaceholderText("PowerPoint (.pptx) ファイルを選択してください")
            browse_button = QPushButton("参照", self)
            browse_button.clicked.connect(self.choose_input_file)

            form_layout.addWidget(self.input_edit)
            form_layout.addWidget(browse_button)

            output_layout = QHBoxLayout()
            layout.addLayout(output_layout)

            self.output_edit = QLineEdit(self)
            self.output_edit.setPlaceholderText("出力フォルダを選択してください")
            output_button = QPushButton("参照", self)
            output_button.clicked.connect(self.choose_output_dir)

            output_layout.addWidget(self.output_edit)
            output_layout.addWidget(output_button)

            self.convert_button = QPushButton("変換開始", self)
            self.convert_button.clicked.connect(self.start_conversion)
            layout.addWidget(self.convert_button)

            self.log_view = QTextEdit(self)
            self.log_view.setReadOnly(True)
            layout.addWidget(self.log_view)

            status = QStatusBar(self)
            self.setStatusBar(status)
            self.status_bar = status

            self.setCentralWidget(central)

            toolbar = QToolBar("メインツールバー", self)
            toolbar.setMovable(False)

            open_action = QAction("PowerPointを開く", self)
            open_action.triggered.connect(self.choose_input_file)
            toolbar.addAction(open_action)

            output_action = QAction("出力先を開く", self)
            output_action.triggered.connect(self.choose_output_dir)
            toolbar.addAction(output_action)

            self.addToolBar(toolbar)

        def append_log(self, message: str) -> None:
            self.log_view.append(message)
            self.log_view.ensureCursorVisible()
            self.status_bar.showMessage(message, 5000)

        def choose_input_file(self) -> None:
            file_path, _ = QFileDialog.getOpenFileName(
                self,
                "入力 PPTX を選択",
                "",
                "PowerPoint ファイル (*.pptx)",
            )
            if file_path:
                self.input_edit.setText(file_path)
                if not self.output_edit.text():
                    self.output_edit.setText(str(Path(file_path).parent))

        def choose_output_dir(self) -> None:
            directory = QFileDialog.getExistingDirectory(
                self,
                "出力フォルダを選択",
                "",
            )
            if directory:
                self.output_edit.setText(directory)

        def start_conversion(self) -> None:
            input_path = self.input_edit.text().strip()
            if not input_path:
                QMessageBox.warning(self, "入力ファイル", "入力ファイルを選択してください。")
                return
            input_file = Path(input_path)
            if not input_file.exists():
                QMessageBox.critical(self, "入力ファイル", "入力ファイルが見つかりません。")
                return

            output_path_value = self.output_edit.text().strip()
            output_dir = Path(output_path_value) if output_path_value else input_file.parent

            if not output_dir.exists():
                try:
                    output_dir.mkdir(parents=True, exist_ok=True)
                except OSError as exc:
                    QMessageBox.critical(self, "出力フォルダ", f"出力フォルダを作成できません: {exc}")
                    return

            self.log_view.clear()
            self.append_log("変換を開始します...")
            self.convert_button.setEnabled(False)
            self.status_bar.showMessage("変換中...", 0)

            self.worker = ConversionThread(input_file, output_dir)
            self.worker.progress.connect(self.append_log)
            self.worker.finished.connect(self.finish_conversion)
            self.worker.start()

        def finish_conversion(self, success: bool, message: str) -> None:
            self.append_log(message)
            self.status_bar.showMessage("完了" if success else "エラー", 5000)
            self.convert_button.setEnabled(True)


    def run_app() -> None:
        app = QApplication(sys.argv)
        window = MainWindow()
        window.show()
        sys.exit(app.exec())


else:

    def run_app() -> None:
        raise RuntimeError(
            "PySide6 またはその依存関係が利用できないため GUI を起動できません。"
            f" 原因: {PYSIDE_IMPORT_ERROR}"
        )


if __name__ == "__main__":
    if not PYSIDE_AVAILABLE:
        raise SystemExit(
            "PySide6 が利用できないため GUI モードを起動できません。"
            f" 原因: {PYSIDE_IMPORT_ERROR}"
        )
    run_app()
